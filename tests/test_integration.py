from pathlib import Path

import pytest
from PIL import Image
from pptx import Presentation

from tiff_to_powerpoint.generator import generate_presentation
from tiff_to_powerpoint.models import AppConfig
from tiff_to_powerpoint.pairing import pair_images
from tiff_to_powerpoint.scanner import scan_folder


def _make_tiff(path: Path, colour: tuple[int, int, int]) -> None:
    Image.new("RGB", (120, 60), colour).save(path, format="TIFF")


def test_end_to_end_generation_has_expected_slide_count_and_aspect_ratio(tmp_path: Path):
    for primary, count in (("E9", 12), ("E8", 3)):
        for sub_component in range(1, count + 1):
            _make_tiff(tmp_path / f"{primary}_{sub_component}.tif", (20, 80, 160))
            _make_tiff(
                tmp_path / f"{primary}_{sub_component}_Volume_Viewer_1.tif",
                (160, 80, 20),
            )

    scan = scan_folder(tmp_path)
    pairing = pair_images(scan.parsed_images)
    output = tmp_path / "result.pptx"
    config = AppConfig(
        root_folder=tmp_path,
        image_width_cm=2.0,
        horizontal_gap_cm=0.2,
        vertical_gap_cm=0.2,
        max_columns_per_slide=5,
        output_path=output,
    )

    result = generate_presentation(pairing, config)

    assert result.output_path == output.resolve()
    assert result.slide_count == 4
    assert output.is_file()
    presentation = Presentation(output)
    assert len(presentation.slides) == 4
    titles = [slide.shapes[0].text for slide in presentation.slides]
    assert titles == ["E8", "E9 - 1/3", "E9 - 2/3", "E9 - 3/3"]
    first_picture = next(shape for shape in presentation.slides[0].shapes if shape.shape_type == 13)
    assert first_picture.width / first_picture.height == pytest.approx(2.0, rel=0.01)
